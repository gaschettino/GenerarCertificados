import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import tempfile
import subprocess
import shutil

# --- Configuración de la página ---
st.set_page_config(page_title="Generador de Certificados", layout="centered")

st.title("Generador de Certificados")
st.write("")
st.write("Cualquier consulta enviar mail a gaschettino@garrahan.gov.ar")

# --- Subida de archivos ---
uploaded_template = st.file_uploader("Subí el template del certificado (.pptx)", type=["pptx"])
uploaded_excel = st.file_uploader("Subí el listado de asistentes (.xlsx). Tiene que tener dos columnas: 'Nombre' y 'Apellido'")

# --- Función para convertir PPTX → PDF usando LibreOffice ---
def convert_to_pdf(input_pptx, output_dir):
    """Convierte un archivo .pptx a .pdf usando LibreOffice (modo headless)."""
    try:
        result = subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_pptx
        ], check=True, capture_output=True, text=True)
        
        # Si la conversión fue exitosa, eliminar el archivo PPTX
        if result.returncode == 0:
            os.remove(input_pptx)
            return True
        return False
    except Exception as e:
        print(f"Error al convertir {input_pptx}: {e}")
        return False

if uploaded_template and uploaded_excel:
    if st.button("🚀 Generar certificados"):
        with st.spinner("Generando certificados..."):
            with tempfile.TemporaryDirectory() as tmpdir:
                # Guardar archivos subidos
                template_path = os.path.join(tmpdir, "template.pptx")
                with open(template_path, "wb") as f:
                    f.write(uploaded_template.read())

                # Leer el Excel
                df = pd.read_excel(uploaded_excel)

                # Normalizar columnas
                df.columns = df.columns.str.strip().str.title()

                # Si tiene columnas "Apellido" y "Nombre", combinarlas
                if "Apellido" in df.columns and "Nombre" in df.columns:
                    df["Nombre y apellido"] = (
                        df["Apellido"].astype(str).str.strip() + " " + df["Nombre"].astype(str).str.strip()
                    ).str.title()
                elif "Nombre Y Apellido" in df.columns:
                    df["Nombre y apellido"] = df["Nombre Y Apellido"].astype(str).str.title()
                else:
                    st.error("No se encontró una columna de nombre adecuada (Apellido/Nombre o Nombre y Apellido).")
                    st.stop()

                # Si existe una columna Asistió, filtrar por 'SI'
                if "Asistió" in df.columns:
                    df = df[df["Asistió"].astype(str).str.upper() == "SI"]

                # Crear carpeta de salida
                output_dir = os.path.join(tmpdir, "Certificados")
                os.makedirs(output_dir, exist_ok=True)

                # Contadores para el progreso
                total_nombres = len(df["Nombre y apellido"].dropna().unique())
                progress_bar = st.progress(0)
                status_text = st.empty()

                # Generar certificados
                for i, nombre in enumerate(df["Nombre y apellido"].dropna().unique()):
                    status_text.text(f"Procesando: {nombre} ({i+1}/{total_nombres})")
                    
                    prs = Presentation(template_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if "Nombre y apellido" in run.text:
                                            run.text = run.text.replace("Nombre y apellido", nombre)
                                            run.font.name = "Quintessential"
                                            run.font.size = Pt(25)
                                            run.font.italic = True
                                            run.font.weight = bold
                                            run.font.color.rgb = RGBColor(0, 0, 0)
                                    paragraph.alignment = PP_ALIGN.CENTER

                    # Guardar PPTX temporal
                    safe_name = "".join(c for c in nombre if c.isalnum() or c in (' ', '-', '_')).rstrip()
                    safe_name = safe_name.replace(' ', '_')
                    output_pptx = os.path.join(output_dir, f"Certificado_{safe_name}.pptx")
                    prs.save(output_pptx)

                    # Convertir a PDF y eliminar PPTX
                    convert_to_pdf(output_pptx, output_dir)
                    
                    # Actualizar barra de progreso
                    progress_bar.progress((i + 1) / total_nombres)

                status_text.text("Procesamiento completado")

                # Verificar que solo hay PDFs en la carpeta
                archivos_finales = os.listdir(output_dir)
                pdf_count = len([f for f in archivos_finales if f.endswith('.pdf')])
                pptx_count = len([f for f in archivos_finales if f.endswith('.pptx')])
                
                st.info(f"Se generaron {pdf_count} archivos PDF. Archivos PPTX eliminados: {pptx_count}")

                # Crear archivo ZIP con todos los certificados PDF
                zip_path = os.path.join(tmpdir, "certificados.zip")
                shutil.make_archive(zip_path.replace(".zip", ""), "zip", output_dir)

                st.success("Certificados PDF generados correctamente.")
                st.write("Podés descargar todos los certificados PDF a continuación:")

                with open(zip_path, "rb") as f:
                    st.download_button(
                        "Descargar Certificados PDF", 
                        f, 
                        "certificados.pdf.zip", 
                        "application/zip"
                    )


